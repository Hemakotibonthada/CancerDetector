"""Generate CancerGuard AI PowerPoint Presentation"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BLUE = RGBColor(0x15, 0x65, 0xC0)
DARK_BLUE = RGBColor(0x0D, 0x47, 0xA1)
PURPLE = RGBColor(0x76, 0x4B, 0xA2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x21, 0x21, 0x21)
GRAY = RGBColor(0x75, 0x75, 0x75)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED = RGBColor(0xC6, 0x28, 0x28)
ORANGE = RGBColor(0xE6, 0x51, 0x00)
ACCENT_BLUE = RGBColor(0x42, 0xA5, 0xF5)

def add_gradient_bg(slide, c1=DARK_BLUE, c2=PURPLE):
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = c1
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = c2
    fill.gradient_stops[1].position = 1.0

def add_solid_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=DARK):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(6)
    return txBox

def add_stat_card(slide, left, top, width, height, number, label, bg_color=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Number
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    p.space_before = Pt(20)
    # Label
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(14)
    p2.font.color.rgb = WHITE
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER

def add_section_header(slide, text, top=0.3):
    add_text_box(slide, 0.8, top, 11, 0.7, text, font_size=32, bold=True, color=BLUE)
    # underline bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(top + 0.65), Inches(2), Inches(0.05))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()

def add_table(slide, left, top, width, rows_data, col_widths=None):
    rows = len(rows_data)
    cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(0.4 * rows))
    table = table_shape.table
    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.name = "Calibri"
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = BLUE
                else:
                    paragraph.font.color.rgb = DARK
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else LIGHT_BG
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

# ============================================================================
# SLIDE 1: Title
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_gradient_bg(slide, DARK_BLUE, PURPLE)
add_text_box(slide, 1, 1.2, 11, 1.2, "CancerGuard AI", font_size=54, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 2.5, 11, 0.8, "AI-Powered Cancer Detection & Comprehensive Healthcare Platform", font_size=24, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 3.5, 11, 0.5, "Version 1.0.0  |  Full-Stack Enterprise Healthcare Solution", font_size=16, color=RGBColor(0xBB, 0xDE, 0xFB), alignment=PP_ALIGN.CENTER)
# stat boxes on title
for i, (num, lbl) in enumerate([("457", "API Endpoints"), ("239", "Database Tables"), ("89", "Frontend Pages"), ("38", "API Modules")]):
    add_stat_card(slide, 1.5 + i * 2.8, 4.6, 2.4, 1.5, num, lbl, [BLUE, RGBColor(0xF5, 0x57, 0x6C), ACCENT_BLUE, GREEN][i])
add_text_box(slide, 1, 6.5, 11, 0.4, "React  •  FastAPI  •  TensorFlow  •  PyTorch  •  React Native", font_size=14, color=RGBColor(0xBB, 0xDE, 0xFB), alignment=PP_ALIGN.CENTER)

# ============================================================================
# SLIDE 2: Problem Statement
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)
add_section_header(slide, "The Problem")
problems = [
    "Cancer is the 2nd leading cause of death globally — 10M+ deaths/year",
    "Late-stage detection reduces 5-year survival rate to as low as 10%",
    "Fragmented health records make it difficult to track risk over time",
    "Patients lack tools to actively manage their own health data",
    "Hospitals need integrated systems for end-to-end clinical care",
]
for i, prob in enumerate(problems):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.3 + i * 0.75), Inches(7), Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFF, 0xEB, 0xEE)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].text = f"  {prob}"
    tf.paragraphs[0].font.size = Pt(15)
    tf.paragraphs[0].font.color.rgb = RED
    tf.paragraphs[0].font.name = "Calibri"

# Solution box
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.3), Inches(11.5), Inches(1.2))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
shape.line.fill.background()
tf = shape.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Our Solution"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = GREEN
p.font.name = "Calibri"
p2 = tf.add_paragraph()
p2.text = "An AI-powered platform combining early cancer risk prediction with comprehensive health monitoring, document management, and hospital operations — accessible to everyone."
p2.font.size = Pt(14)
p2.font.color.rgb = DARK
p2.font.name = "Calibri"

# ============================================================================
# SLIDE 3: Technology Stack
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)
add_section_header(slide, "Technology Stack")

tech_data = [
    ("Backend", "FastAPI 0.104\nPython 3.13\nSQLAlchemy 2.0 Async\nJWT Authentication", BLUE),
    ("Frontend", "React 18\nTypeScript 4.9\nMaterial UI v5\nRecharts / Chart.js", RGBColor(0x7B, 0x1F, 0xA2)),
    ("Mobile", "Expo 50\nReact Native 0.73\nReact Navigation 6\nReact Native Paper", GREEN),
    ("AI / ML", "TensorFlow 2.15\nPyTorch 2.1\nXGBoost / LightGBM\nSHAP / LIME", ORANGE),
    ("Database", "SQLite (dev)\nPostgreSQL (prod)\nAlembic Migrations\nAsync ORM", RGBColor(0x00, 0x69, 0x5C)),
    ("Security", "JWT + Refresh Tokens\nbcrypt Hashing\nRBAC (4 roles)\nAudit Logging", RED),
]
for i, (title, details, color) in enumerate(tech_data):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 4.1
    y = 1.3 + row * 2.9
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.7), Inches(2.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.color.rgb = color
    shape.line.width = Pt(3)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.space_after = Pt(8)
    for line in details.split("\n"):
        p2 = tf.add_paragraph()
        p2.text = f"  {line}"
        p2.font.size = Pt(13)
        p2.font.color.rgb = GRAY
        p2.font.name = "Calibri"
        p2.space_after = Pt(2)

# ============================================================================
# SLIDE 4: Three Portals
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)
add_section_header(slide, "Three Dedicated Portals")

portals = [
    ("User Portal", "37", "Pages", [
        "Dashboard & Health Score", "AI Cancer Risk Assessment",
        "Health Records & Blood Tests", "Document Upload & Management",
        "Insurance & Billing", "Telehealth & Messaging",
        "Wearables & Smartwatch", "Mental Health & Wellness",
        "Genomics & Genetic Profile", "Nutrition & Exercise",
    ], BLUE),
    ("Hospital Portal", "30", "Pages", [
        "Operations Dashboard", "Patient Management",
        "Lab & Radiology (AI)", "Pathology & Specimens",
        "Pharmacy & Formulary", "Emergency & Triage",
        "Surgery Scheduling", "Telemedicine",
        "Clinical Trials", "Quality & Safety",
    ], RGBColor(0xF5, 0x57, 0x6C)),
    ("Admin Portal", "22", "Pages", [
        "System Monitoring", "User Management",
        "Hospital Management", "AI Model Management",
        "Audit Logs", "Security & Compliance",
        "Billing & Revenue", "Platform Analytics",
        "Configuration", "Workforce Management",
    ], GREEN),
]
for i, (title, count, unit, features, color) in enumerate(portals):
    x = 0.6 + i * 4.2
    # Header card
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.2), Inches(3.8), Inches(1.3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = count
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    p2 = tf.add_paragraph()
    p2.text = f"{title}"
    p2.font.size = Pt(16)
    p2.font.color.rgb = WHITE
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER
    # Features
    add_bullet_list(slide, x + 0.1, 2.7, 3.6, 4.5, [f"• {f}" for f in features], font_size=12, color=GRAY)

# ============================================================================
# SLIDE 5: User Portal Features
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)
add_section_header(slide, "User Portal — Key Features")

features_left = [
    "🎯  AI Cancer Risk Assessment (multi-model ensemble)",
    "📊  Health Dashboard with live vital signs",
    "📁  Medical Document Upload & Management",
    "🛡️  Insurance Policy & Claims Tracking",
    "⌚  Smartwatch & Wearable Integration",
    "🩸  Blood Tests & Biomarker Tracking",
    "💊  Medication Management & Adherence",
    "📅  Appointment Scheduling & Reminders",
    "🧬  Genetic Profile & Genomics",
    "🧠  Mental Health (CBT, Mindfulness, Crisis)",
]
features_right = [
    "🥗  Diet, Nutrition & Hydration Logging",
    "🏃  Exercise & Fitness Tracking",
    "📹  Telehealth Video Consultations",
    "💬  Secure Messaging with Care Team",
    "📚  Health Education & Quizzes",
    "🩺  Screening Schedule & Reminders",
    "👨‍👩‍👧  Family Health History & Hereditary Risk",
    "🩸  Blood Donor Registration & Matching",
    "💰  Billing & Payment Tracking",
    "🏆  Gamification & Health Challenges",
]
add_bullet_list(slide, 0.6, 1.2, 5.8, 6, features_left, font_size=14, color=DARK)
add_bullet_list(slide, 6.8, 1.2, 5.8, 6, features_right, font_size=14, color=DARK)

# ============================================================================
# SLIDE 6: Hospital Portal Features
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)
add_section_header(slide, "Hospital Portal — Clinical Operations")

hosp_left = [
    "📋  Patient Management & Demographics",
    "🏥  Bed Management & Occupancy Tracking",
    "🔬  Lab Management & Quality Control",
    "📷  Radiology with AI-Assisted Reads",
    "🧫  Pathology (Specimens, Slides, Tumor Board)",
    "💊  Pharmacy & Formulary Management",
    "🚑  Emergency (Triage, Sepsis, Stroke Screening)",
    "🩸  Blood Bank & Cross-Matching",
    "🔪  Surgery Scheduling & Case Tracking",
]
hosp_right = [
    "📹  Telemedicine & Virtual Waiting Room",
    "🧬  Genomics Lab Workflows",
    "📊  Clinical Decision Support & Guidelines",
    "🧪  Clinical Trials Management",
    "✅  Quality & Safety (Incidents, Infection Control)",
    "📦  Supply Chain & Asset Tracking",
    "👥  Staff Directory & Workforce Management",
    "📊  Population Health & Disease Registries",
    "📈  AI Analytics & Predictive Models",
]
add_bullet_list(slide, 0.6, 1.2, 5.8, 6, hosp_left, font_size=14, color=DARK)
add_bullet_list(slide, 6.8, 1.2, 5.8, 6, hosp_right, font_size=14, color=DARK)

# ============================================================================
# SLIDE 7: System Architecture
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)
add_section_header(slide, "System Architecture")

layers = [
    ("Client Layer", "React Web App (89 pages)  |  Expo Mobile App  |  Swagger API Docs", RGBColor(0xE3, 0xF2, 0xFD), BLUE, 1.2),
    ("", "⬇  HTTPS / REST API  ⬇", WHITE, GRAY, 2.2),
    ("API Gateway", "CORS Middleware → JWT Authentication → Request Timing → Audit Logging", RGBColor(0xE8, 0xF5, 0xE9), GREEN, 2.7),
    ("Application Layer", "38 API Router Modules  →  Service Layer  →  Pydantic Validation  |  457 Endpoints", RGBColor(0xFD, 0xF5, 0xE6), ORANGE, 3.6),
    ("Data Access Layer", "SQLAlchemy 2.0 Async ORM  |  40 Model Files  |  239 Tables  |  UUID PKs  |  Soft Delete", RGBColor(0xF3, 0xE5, 0xF5), PURPLE, 4.5),
    ("Database", "SQLite (development)  |  PostgreSQL (production)  |  Alembic Migrations", LIGHT_BG, GRAY, 5.4),
    ("AI/ML Engine", "TensorFlow  |  PyTorch  |  XGBoost  |  LightGBM  |  CatBoost  |  SHAP / LIME", RGBColor(0xFF, 0xF3, 0xE0), ORANGE, 6.1),
]
for title, desc, bg, fg, y in layers:
    if title:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(y), Inches(11.3), Inches(0.7))
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg
        shape.line.color.rgb = fg
        shape.line.width = Pt(1.5)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"{title}:  {desc}"
        p.font.size = Pt(13)
        p.font.color.rgb = fg
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
    else:
        add_text_box(slide, 1, y, 11.3, 0.4, desc, font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================================
# SLIDE 8: Data Model
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)
add_section_header(slide, "Data Model — 239 Tables")

table_data = [
    ["Domain", "Tables", "Key Models"],
    ["User & Auth", "4", "User, Session, Preference"],
    ["Patient & Health", "14", "Patient, Records, Vitals, Labs"],
    ["Cancer & Screening", "6", "Screening, RiskAssessment, Prediction"],
    ["Wearables & IoT", "22", "Glucose, ECG, Gait, Sleep, SpO2"],
    ["Genomics", "8", "Sequencing, Variants, LiquidBiopsy"],
    ["Clinical Trials", "8", "Protocol, Participants, Visits"],
    ["Imaging & Pathology", "17", "Radiology, AI Reads, Specimens"],
    ["Billing & Insurance", "16", "Invoice, Claims, PriorAuth"],
    ["Communication", "15", "Messages, Referrals, Telehealth"],
    ["Quality & Safety", "9", "Incidents, Infections, Checklists"],
    ["Mental Health", "9", "CBT, Mindfulness, Crisis"],
    ["Nutrition & Rehab", "15", "MealPlan, FoodLog, TherapySessions"],
    ["Education & Engagement", "21", "Resources, Quizzes, Gamification"],
    ["Supply Chain & Workforce", "13", "Inventory, Equipment, Scheduling"],
    ["Social & Emergency", "13", "SDOH, Triage, Sepsis, Trauma"],
    ["System & Documents", "9", "Config, Audit, Documents, Notifications"],
]
add_table(slide, 0.8, 1.2, 11.5, table_data)

# ============================================================================
# SLIDE 9: Security
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)
add_section_header(slide, "Security Architecture")

security_layers = [
    ("Layer 1: CORS", "Origin-based access control whitelist", RGBColor(0xE3, 0xF2, 0xFD)),
    ("Layer 2: JWT Auth", "Access tokens (15 min) + Refresh tokens (7 days)", RGBColor(0xE8, 0xF5, 0xE9)),
    ("Layer 3: RBAC", "Role-Based Access: Patient / Doctor / Hospital Admin / System Admin", RGBColor(0xFD, 0xF5, 0xE6)),
    ("Layer 4: Validation", "Pydantic input validation on all API requests", RGBColor(0xFC, 0xE4, 0xEC)),
    ("Layer 5: ORM Protection", "SQLAlchemy parameterized queries — SQL injection prevention", RGBColor(0xF3, 0xE5, 0xF5)),
    ("Layer 6: Password Security", "bcrypt hashing with configurable work factor", RGBColor(0xE0, 0xF2, 0xF1)),
    ("Layer 7: Audit Trail", "Comprehensive logging of all critical operations", RGBColor(0xFB, 0xF0, 0xE3)),
    ("Layer 8: Soft Delete", "Data preservation — records never permanently deleted", LIGHT_BG),
]
for i, (title, desc, bg) in enumerate(security_layers):
    y = 1.2 + i * 0.72
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(y), Inches(10), Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"  🔒  {title}  —  {desc}"
    p.font.size = Pt(14)
    p.font.color.rgb = DARK
    p.font.name = "Calibri"

# ============================================================================
# SLIDE 10: AI/ML Pipeline
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)
add_section_header(slide, "AI/ML Cancer Risk Engine")

pipeline_steps = [
    ("1", "Data\nCollection", "Blood biomarkers\nGenetics\nLifestyle factors"),
    ("2", "Pre-\nprocessing", "Feature engineering\nNormalization\nMissing data"),
    ("3", "Ensemble\nModels", "TensorFlow DNN\nPyTorch\nXGBoost/LightGBM"),
    ("4", "Risk\nPrediction", "Risk score (0-1)\nRisk level\nConfidence"),
    ("5", "Explain-\nability", "SHAP values\nLIME\nFeature importance"),
]
for i, (num, title, desc) in enumerate(pipeline_steps):
    x = 0.6 + i * 2.55
    # Circle
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.7), Inches(1.4), Inches(0.7), Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    # Title
    add_text_box(slide, x, 2.3, 2.2, 0.8, title, font_size=14, bold=True, color=BLUE, alignment=PP_ALIGN.CENTER)
    # Description
    add_text_box(slide, x, 3.1, 2.2, 1.2, desc, font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

# Risk levels bar
levels = [
    ("Very Low", RGBColor(0x4C, 0xAF, 0x50)),
    ("Low", RGBColor(0x8B, 0xC3, 0x4A)),
    ("Moderate", RGBColor(0xFF, 0x98, 0x00)),
    ("High", RGBColor(0xF4, 0x43, 0x36)),
    ("Very High", RGBColor(0xD3, 0x2F, 0x2F)),
    ("Critical", RGBColor(0xB7, 0x1C, 0x1C)),
]
add_text_box(slide, 0.8, 4.6, 4, 0.5, "Risk Levels:", font_size=16, bold=True, color=DARK)
for i, (lbl, clr) in enumerate(levels):
    x = 0.8 + i * 2
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(5.1), Inches(1.8), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = clr
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = lbl
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

# Factors
add_text_box(slide, 0.8, 5.9, 12, 0.5, "Input Factors:  Blood biomarkers  •  Genetic variants  •  Family history  •  Lifestyle  •  Vitals  •  Medical history", font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================================
# SLIDE 11: Key Differentiators
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)
add_section_header(slide, "Key Differentiators")

diffs = [
    ("🧠", "Multi-Model AI Ensemble", "Not just one algorithm — combines deep learning (TensorFlow, PyTorch), gradient boosting (XGBoost, LightGBM, CatBoost), and classical ML for robust, explainable predictions."),
    ("📱", "Cross-Platform", "Consistent experience across Web (React), iOS, and Android (Expo/React Native) with a single shared FastAPI backend."),
    ("🏥", "Full Hospital EHR", "Not just patient-facing — comprehensive clinical operations: radiology, pathology, pharmacy, surgery, emergency, genomics lab."),
    ("⌚", "IoT & Wearable Integration", "22 wearable data tables for continuous monitoring: glucose (CGM), ECG, gait analysis, respiratory, sleep, fall detection."),
    ("🔬", "Genomics-Ready", "Built-in support for liquid biopsy, gene panels, pharmacogenomics, hereditary cancer assessment, and gene expression analysis."),
]
for i, (icon, title, desc) in enumerate(diffs):
    y = 1.2 + i * 1.15
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y), Inches(11.5), Inches(1.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{icon}  {title}"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p.font.name = "Calibri"
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(12)
    p2.font.color.rgb = GRAY
    p2.font.name = "Calibri"

# ============================================================================
# SLIDE 12: API Overview
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)
add_section_header(slide, "API Overview — 457 Endpoints")

api_data = [
    ["Module", "Endpoints", "Module", "Endpoints", "Module", "Endpoints"],
    ["Auth", "6", "Billing", "19", "Education", "16"],
    ["Users", "5", "Quality & Safety", "18", "Social Determinants", "16"],
    ["Patients", "7", "Communication", "19", "Supply Chain", "16"],
    ["Hospitals", "7", "Telehealth", "17", "Wearables", "22"],
    ["Cancer Detection", "3", "Clinical Trials", "17", "Emergency", "14"],
    ["Health Records", "4", "Pharmacy", "17", "Workforce", "15"],
    ["Blood Samples", "5", "Pathology", "16", "Documents", "14"],
    ["Appointments", "3", "Mental Health", "18", "Blood Donor", "14"],
    ["Genomics", "15", "Nutrition", "16", "Patient Engagement", "21"],
    ["Clinical Decision", "16", "Rehabilitation", "15", "Research", "15"],
    ["Smartwatch", "4", "Population Health", "16", "Admin", "5"],
    ["Radiology", "13", "Notifications", "3", "Analytics", "2"],
]
add_table(slide, 0.5, 1.2, 12.3, api_data)

# ============================================================================
# SLIDE 13: Demo Accounts
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide, WHITE)
add_section_header(slide, "Demo & Access")

demo_data = [
    ["Role", "Email", "Password"],
    ["👤  User (Patient)", "patient@cancerguard.ai", "Patient@123456"],
    ["🩺  Doctor", "doctor@cancerguard.ai", "Doctor@123456"],
    ["🏥  Hospital Admin", "hospital.admin@cancerguard.ai", "Hospital@123456"],
    ["⚙️  System Admin", "admin@cancerguard.ai", "Admin@123456"],
]
add_table(slide, 2, 1.4, 9, demo_data, col_widths=[2.5, 3.5, 3])

urls = [
    ("Backend API", "http://localhost:8000"),
    ("Swagger Docs", "http://localhost:8000/docs"),
    ("Frontend App", "http://localhost:3000"),
    ("Mobile App", "Expo Go — scan QR code"),
    ("GitHub", "github.com/Hemakotibonthada/CancerDetector"),
]
for i, (lbl, url) in enumerate(urls):
    add_text_box(slide, 2, 4.2 + i * 0.5, 3, 0.4, lbl, font_size=14, bold=True, color=BLUE)
    add_text_box(slide, 5, 4.2 + i * 0.5, 7, 0.4, url, font_size=14, color=GRAY)

# ============================================================================
# SLIDE 14: Thank You
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide, DARK_BLUE, PURPLE)
add_text_box(slide, 1, 2.0, 11, 1.2, "Thank You", font_size=54, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 3.3, 11, 0.8, "CancerGuard AI — Empowering Early Detection, Saving Lives", font_size=22, color=RGBColor(0xBB, 0xDE, 0xFB), alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 4.8, 11, 0.5, "github.com/Hemakotibonthada/CancerDetector", font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 5.5, 11, 0.4, "Built with  FastAPI  •  React  •  TensorFlow  •  Expo", font_size=14, color=RGBColor(0xBB, 0xDE, 0xFB), alignment=PP_ALIGN.CENTER)

# ============================================================================
# Save
# ============================================================================
output_path = r"c:\Users\v-hbonthada\WorkSpace-Pract\Cancer detection\docs\CancerGuard_AI_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
